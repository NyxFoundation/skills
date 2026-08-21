#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0"]
# ///
"""Slidev の PNG 書き出しを .pptx に包む。

    uv run scripts/png-to-pptx.py dist-png out.pptx [notes.md]

会場が .pptx / Google Slides しか受け付けないときの逃げ道。
1 スライド = 全面 1 枚の画像なので、Slidev 側のデザイン（webfont・SVG・レイアウト）が
そのまま保存される。文字は選択・編集できなくなるが、崩れることもない。

notes.md を渡すと、スピーカーノートを流し込む。書式は `---` 区切りで 1 枚 1 ブロック：

    01 枚目の口頭メモ
    ---
    02 枚目の口頭メモ

前提: `slidev export --format png --per-slide` の出力（NN.png）。
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# 16:9 / 13.333in × 7.5in
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)


def natural_key(p: Path):
    m = re.search(r"(\d+)", p.stem)
    return (int(m.group(1)) if m else 0, p.stem)


def main() -> None:
    if len(sys.argv) < 3:
        sys.exit(__doc__)
    png_dir = Path(sys.argv[1])
    out = Path(sys.argv[2])
    notes_path = Path(sys.argv[3]) if len(sys.argv) > 3 else None

    pngs = sorted(png_dir.glob("*.png"), key=natural_key)
    if not pngs:
        sys.exit(f"PNG が見つかりません: {png_dir}")

    notes = []
    if notes_path and notes_path.exists():
        notes = [b.strip() for b in notes_path.read_text(encoding="utf-8").split("\n---\n")]
        if len(notes) != len(pngs):
            print(f"警告: ノート {len(notes)} 件 / スライド {len(pngs)} 枚。数が合いません。",
                  file=sys.stderr)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for i, png in enumerate(pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, SLIDE_W, SLIDE_H)
        if i < len(notes) and notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    size_mb = out.stat().st_size / 1024 / 1024
    print(f"{len(pngs)} 枚 → {out} ({size_mb:.1f} MB)")
    if size_mb > 40:
        print("注意: 40MB 超。メール添付には向きません。", file=sys.stderr)


if __name__ == "__main__":
    main()
