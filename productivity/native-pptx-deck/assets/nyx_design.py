"""ネイティブ .pptx のデザイン層（参考実装）。

内容側のコードから OOXML を一切見せないための薄いレイヤ。
トークンと書体名を差し替えれば別プロジェクトでそのまま使える。

公式 API が無いもの（letter-spacing / 塗りの透明度 / 曲線のスムージング / 影 /
group transform）は、ここで OOXML を直接叩いて隠している。
詳細は references/ooxml-workarounds.md。
"""

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ── トークン（HP tokens.css 準拠・固定） ──────────────────────────────
BG        = RGBColor(0xFA, 0xF9, 0xF5)
BG2       = RGBColor(0xF3, 0xF1, 0xEA)
INK       = RGBColor(0x18, 0x18, 0x1A)
INK_DIM   = RGBColor(0x55, 0x52, 0x4C)
INK_FAINT = RGBColor(0x9A, 0x95, 0x8C)
LINE      = RGBColor(0xE3, 0xE1, 0xDB)   # rgba(24,24,26,.10) を bg 上でフラット化
LINE_STR  = RGBColor(0xCF, 0xCD, 0xC6)
ACCENT    = RGBColor(0x1F, 0x3A, 0x52)
SEVERE    = RGBColor(0xA2, 0x54, 0x34)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ── 書体 ────────────────────────────────────────────────────────────
SERIF    = "Cormorant Garamond"
JP       = "Shippori Mincho"
MONO     = "JetBrains Mono"
WORDMARK = "BIZ UDPMincho"

# ── 16:9 キャンバス ─────────────────────────────────────────────────
W, H = Emu(12192000), Emu(6858000)
MARGIN = Emu(838200)          # 0.9in ≒ style.css の 2.6rem 相当

def inches(v):
    return Emu(int(v * 914400))


CONTENT_W = W - MARGIN * 2

# ── グリッド（12 カラム × 8pt ベースライン） ────────────────────────
# 内容側では inches(2.35) のような数字を書かず、col() と band() だけを使う。
# 座標をその場で決めると崩れが再発する（references/layout-and-text.md）。
COLS = 12
GUTTER = Emu(114300)          # 0.125in
BASE = Emu(101600)            # 8pt


def col(start: int, span: int = 1):
    """1-indexed のカラム位置。(x, width) を返す。"""
    unit = (CONTENT_W - GUTTER * (COLS - 1)) / COLS
    x = MARGIN + (unit + GUTTER) * (start - 1)
    return Emu(int(x)), Emu(int(unit * span + GUTTER * (span - 1)))


def band(n: float):
    """8pt を 1 バンドとした縦位置。"""
    return Emu(int(BASE * n))


# ── OOXML を直接叩くヘルパ（公式 API が無いもの） ────────────────────

def set_spacing(run, em):
    """letter-spacing。python-pptx に API が無いので rPr@spc を直接書く。
    em はフォントサイズに対する比（キッカーは 0.18）。単位は 1/100pt。"""
    size_pt = run.font.size.pt if run.font.size else 18
    run.font._rPr.set("spc", str(int(size_pt * em * 100)))


def set_alpha(spPr, alpha_pct: int):
    """塗り／線の透明度。a:srgbClr の下に a:alpha を差し込む。"""
    srgb = spPr.find(".//" + qn("a:srgbClr"))
    if srgb is None:
        return
    a = srgb.makeelement(qn("a:alpha"), {"val": str(alpha_pct * 1000)})
    srgb.append(a)


def set_smooth(chart, on=True):
    """曲線のスムージング。XySeries に smooth プロパティが存在せず、
    `ser.smooth = True` は黙って無視される（属性が生えるだけ）。
    ここで c:smooth を直接書き換える。"""
    for s in chart._chartSpace.iter(qn("c:smooth")):
        s.set("val", "1" if on else "0")


def noshadow(shape):
    """影を消す。

    `shadow.inherit = False` だけでは足りない。add_shape / add_connector が付ける
    `<p:style>` の `effectRef` からテーマの影を継承してしまうため、
    spPr に空の `a:effectLst` を入れたうえで `<p:style>` ごと外す。
    テーマ既定の影は禁則（線は 1px の軽量フレームで引く）。"""
    el = shape._element
    spPr = el.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    for style in el.findall(qn("p:style")):
        el.remove(style)
    return shape


_BODYPR_ORDER = ("a:prstTxWarp", "a:noAutofit", "a:normAutofit", "a:spAutoFit",
                 "a:scene3d", "a:sp3d", "a:flatTx", "a:extLst")


def no_autofit(tf):
    """テキストの自動縮小を切る（サイズ階層を守るため）。

    bodyPr の子要素はスキーマ順が厳格なので、末尾に append せず正しい位置に挿す。"""
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    el = bodyPr.makeelement(qn("a:noAutofit"), {})
    after = _BODYPR_ORDER.index("a:noAutofit")
    pos = len(bodyPr)
    for i, child in enumerate(bodyPr):
        tag = child.tag.split("}")[-1]
        full = "a:" + tag
        if full in _BODYPR_ORDER and _BODYPR_ORDER.index(full) > after:
            pos = i
            break
    bodyPr.insert(pos, el)


# ── テキスト計測（python-pptx に無いので近似する） ──────────────────

def _em_width(text: str) -> float:
    """文字列の幅を em 単位で概算する。
    全角（CJK・全角記号）は 1.0em、半角は 0.5em として数える。
    厳密なメトリクスは取れないので、余白を見込んで使うこと。"""
    w = 0.0
    for c in text:
        o = ord(c)
        full = (0x3000 <= o <= 0x30FF or 0x4E00 <= o <= 0x9FFF
                or 0xFF00 <= o <= 0xFF60 or 0xFFE0 <= o <= 0xFFE6
                or 0x3400 <= o <= 0x4DBF)
        w += 1.0 if full else 0.5
    return w


def est_lines(text: str, width_emu, size_pt: float) -> int:
    """折り返し後の行数を概算する。"""
    if not text:
        return 1
    width_pt = width_emu / 12700.0
    per_line = max(1.0, width_pt / size_pt)
    import math
    return max(1, math.ceil(_em_width(text) / per_line))


def est_height(lines, width_emu, size_pt: float, line_ratio: float = 1.8,
               safety: float = 1.18):
    """body() に渡す行のリストから、必要な高さ（EMU）を概算する。

    実測メトリクスが取れないので必ず外す。safety で余裕を持たせておかないと、
    1 行ぶん溢れてカードの外に文字が出る（PowerPoint と LibreOffice で
    折り返し位置も違うので、詰めすぎない）。"""
    total = 0
    for ln in lines:
        text = ln if isinstance(ln, str) else "".join(
            (sg if isinstance(sg, str) else sg[0]) for sg in ln)
        total += est_lines(text, width_emu, size_pt)
    return Emu(int(total * size_pt * line_ratio * 12700 * safety))


# ── 描画プリミティブ ────────────────────────────────────────────────

def fix_group_xfrm(slide):
    """spTree の `<p:grpSpPr/>` が空だと group の transform が null になり、
    厳しめのインポータが落ちる。明示的に単位変換を書き込む。"""
    spTree = slide.shapes._spTree
    grpSpPr = spTree.find(qn("p:grpSpPr"))
    if grpSpPr is None:
        return
    if grpSpPr.find(qn("a:xfrm")) is not None:
        return
    xfrm = grpSpPr.makeelement(qn("a:xfrm"), {})
    for tag, attrs in (("a:off", {"x": "0", "y": "0"}),
                       ("a:ext", {"cx": str(int(W)), "cy": str(int(H))}),
                       ("a:chOff", {"x": "0", "y": "0"}),
                       ("a:chExt", {"cx": str(int(W)), "cy": str(int(H))})):
        xfrm.append(xfrm.makeelement(qn(tag), attrs))
    grpSpPr.insert(0, xfrm)


def canvas(slide):
    """スライド全面を bg で塗る。"""
    from pptx.enum.shapes import MSO_SHAPE
    fix_group_xfrm(slide)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG
    sh.line.fill.background()
    noshadow(sh)
    return sh


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    no_autofit(tf)
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def run(p, text, *, font=JP, size=14, color=INK, italic=False, bold=False,
        spacing=None):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.bold = bold
    # 日本語グリフにも同じ書体を当てる（既定だと ea フォントが別になる）
    rPr = r.font._rPr
    for tag in ("a:ea", "a:cs"):
        el = rPr.makeelement(qn(tag), {"typeface": font})
        rPr.append(el)
    if spacing:
        set_spacing(r, spacing)
    return r


def para(tf, *, space_after=0, line=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        pct = lnSpc.makeelement(qn("a:spcPct"), {"val": str(int(line * 100000))})
        lnSpc.append(pct)
        pPr.insert(0, lnSpc)
    return p


def rule(slide, x, y, w, color=LINE, weight=1.0):
    """細い罫線。黒太枠は禁則なので必ず 1px 相当で引く。

    コネクタ（`add_connector`）は使わない。水平線だと ext の cy が 0 になり、
    退化した transform を受け付けないインポータがある
    （Google Slides は ShapeImporter の transform で NPE を出す）。
    高さを持った薄い矩形で引く。描画結果も安定する。"""
    from pptx.enum.shapes import MSO_SHAPE
    h = Emu(max(9525, int(weight * 12700)))      # 最低 1px 相当
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(int(w)), h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    noshadow(sh)
    return sh


def kicker(slide, x, y, text):
    """罫線 + Mono UPPER + 0.18em。`01 ／ ラグ` 形式。"""
    rule(slide, x, y, inches(1.5), LINE_STR, 1.0)
    tb, tf = textbox(slide, x, y + Emu(76200), inches(4), inches(0.28))
    p = para(tf, first=True)
    run(p, text, font=MONO, size=12, color=INK_DIM, spacing=0.18)
    return tb


def _is_latin(text: str) -> bool:
    """ラテン文字・数字だけで出来ているか。"""
    return all(ord(c) < 0x2E80 for c in text if not c.isspace())


def display(slide, x, y, w, parts, size=36):
    """見出し。parts は (text, is_em) の列。

    em の描き分けに注意：**Cormorant Garamond に日本語グリフが無い**。
    日本語に italic Cormorant を当てると別フォントに落ちて偽斜体になり、
    かなり気持ち悪い見た目になる。
    ラテン文字・数字のときだけ italic Cormorant を使い、
    日本語の強調は accent 色だけで示す（傾けない）。"""
    tb, tf = textbox(slide, x, y, w, inches(1.4))
    p = para(tf, first=True, line=1.18)
    for text, is_em in parts:
        latin = _is_latin(text)
        run(p, text,
            font=SERIF if (is_em and latin) else JP,
            size=size,
            color=ACCENT if is_em else INK,
            italic=bool(is_em and latin))
    return tb


def body(slide, x, y, w, h, lines, *, size=14, color=INK_DIM, line=1.8,
         font=JP, align=PP_ALIGN.LEFT):
    """本文。lines は文字列、または (text, bold) タプルの列のリスト。"""
    def norm(seg):
        """str / (str, bold) のどちらでも受ける。"""
        if isinstance(seg, str):
            return (seg, False)
        return (seg[0], bool(seg[1]) if len(seg) > 1 else False)

    tb, tf = textbox(slide, x, y, w, h)
    for i, ln in enumerate(lines):
        p = para(tf, first=(i == 0), line=line, align=align, space_after=0)
        segs = [norm(sg) for sg in ln] if isinstance(ln, list) else [(ln, False)]
        for text, strong in segs:
            run(p, text, font=font, size=size,
                color=INK if strong else color, bold=False)
    return tb


def card(slide, x, y, w, h, *, highlight=False):
    """bg-2 + line 1px の軽量フレーム。highlight のときだけ accent 2px。"""
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    sh.adjustments[0] = 0.045
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG2 if highlight else WHITE
    sh.line.color.rgb = ACCENT if highlight else LINE_STR
    sh.line.width = Pt(2 if highlight else 1)
    noshadow(sh)
    return sh


def verse(slide, x, y, w, h, text, *, size=15):
    """左 2px 縦罫 + bg-2 背景の 1 行リード。"""
    from pptx.enum.shapes import MSO_SHAPE
    bgsh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    bgsh.fill.solid()
    bgsh.fill.fore_color.rgb = BG2
    bgsh.line.fill.background()
    noshadow(bgsh)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)), Emu(25400), Emu(int(h)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    noshadow(bar)
    tb, tf = textbox(slide, x + inches(0.22), y + inches(0.1), w - inches(0.4), h)
    p = para(tf, first=True, line=1.5)
    run(p, text, font=JP, size=size, color=INK)
    return bgsh
